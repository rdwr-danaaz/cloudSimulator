#!/usr/bin/env python3
"""Generate the SOC-X simulator slide deck using the branded Radware template.
Built on top of "ASN Support in Attack Remediation.pptx" so every slide inherits
the Radware theme and logo: cover/closing use the "Cover - General 1" layout;
content slides use the "2 column Bullets slide" layout (branded background +
top-left Radware logo + title + two content columns).
Run:
    pip install python-pptx
    python deploy/generate_presentation.py

By default this writes to PRESENTATION_generated.pptx so it never overwrites a
hand-edited PRESENTATION.pptx (the curated, final deck). Pass --force to write
directly to PRESENTATION.pptx.
"""
from __future__ import annotations
import os
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "ASN Support in Attack Remediation.pptx")
FINAL = os.path.join(ROOT, "PRESENTATION.pptx")
OUT = os.path.join(ROOT, "PRESENTATION_generated.pptx")
LAYOUT_COVER = 11
LAYOUT_CONTENT = 12
COVER_TITLE, COVER_LINE1, COVER_LINE2, COVER_LINE3 = 15, 16, 17, 18
CONTENT_TITLE, CONTENT_LEFT, CONTENT_RIGHT = 10, 1, 12
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
RADWARE_BLUE = RGBColor(0x00, 0x35, 0x82)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF7)
CODE_BORDER = RGBColor(0xD5, 0xDA, 0xE0)
CONTENT_X = Emu(1508125)
CONTENT_TOP = Emu(1377535)
CONTENT_W = Emu(9566276)
def _prs_from_template():
    prs = Presentation(TEMPLATE)
    xml_slides = prs.slides._sldIdLst
    rns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    for sldId in list(xml_slides):
        prs.part.drop_rel(sldId.get(rns))
        xml_slides.remove(sldId)
    return prs
def _layout(prs, idx):
    return prs.slide_masters[0].slide_layouts[idx]
def _ph(slide, idx):
    return slide.placeholders[idx]
def _remove_ph(slide, idx):
    try:
        ph = slide.placeholders[idx]
    except KeyError:
        return
    ph._element.getparent().remove(ph._element)
def _set_title(slide, text):
    _ph(slide, CONTENT_TITLE).text = text
def _fill_bullets(ph, items):
    tf = ph.text_frame
    tf.clear()
    first = True
    for text, level, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.bold = bold
def _code_shape(slide, text, top, height, left=CONTENT_X, width=CONTENT_W, size=12):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = CODE_BG
    shp.line.color.rgb = CODE_BORDER
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = "Consolas"
        run.font.size = Pt(size)
        run.font.color.rgb = INK
    return shp
def _text_shape(slide, top, height, left=CONTENT_X, width=CONTENT_W):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box.text_frame
def _para(tf, text, size=16, bold=False, color=INK, space_after=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p
def build():
    prs = _prs_from_template()
    cover = _layout(prs, LAYOUT_COVER)
    content = _layout(prs, LAYOUT_CONTENT)
    s = prs.slides.add_slide(cover)
    _ph(s, COVER_TITLE).text = "SOC-X Recommendation Simulator"
    _ph(s, COVER_LINE1).text = "A drop-in mock of the SOC-X cloud recommendation API"
    _ph(s, COVER_LINE2).text = "Point any Cyber Controller (ADE) at it - no real cloud required"
    _ph(s, COVER_LINE3).text = "August 2026"
    s = prs.slides.add_slide(content)
    _set_title(s, "What the project solves")
    _fill_bullets(_ph(s, CONTENT_LEFT), [
        ("The challenge", 0, True),
        ("Testing/demoing the anomaly-detection-engine (ADE) needs the live SOC-X cloud that serves protection recommendations", 1, False),
        ("That cloud is slow to set up, hard to control, and impossible to run offline", 1, False),
        ("The solution", 0, True),
        ("A local simulator you fully control that replaces the cloud recommendation service", 1, False),
    ])
    _fill_bullets(_ph(s, CONTENT_RIGHT), [
        ("What it delivers", 0, True),
        ("Deterministic, pinned rule sets for specific networks", 1, False),
        ("Generates rules for any other network on demand", 1, False),
        ("Edit the recommendation JSON from a web UI", 1, False),
        ("Auto-wires ADE over SSH (config + TLS trust + restart)", 1, False),
        ("Many Cyber Controllers share one simulator at once", 1, False),
        ("Runs anywhere in one command (Docker), over HTTPS", 1, False),
    ])
    s = prs.slides.add_slide(content)
    _set_title(s, "Architecture")
    _remove_ph(s, CONTENT_LEFT)
    _remove_ph(s, CONTENT_RIGHT)
    diagram = (
        "                 +--------------------------------------------+\n"
        "   Cyber         |  Simulator host (Docker)                   |\n"
        " Controller  --->|  main.py  (ASGI: mounts at / and /socx_sim)|\n"
        "   (ADE)   HTTPS |     |                                      |\n"
        " _getRecommend.  |  cloud_mock_server.py  (FastAPI)           |\n"
        "      |          |   |-- /api/.../_getRecommendation          |\n"
        "      |          |   |-- /ui  (+ /ui/* JSON APIs)             |\n"
        "      |          |   |-- permanent_responses.py  (pinned)     |\n"
        "      |          |   +-- response_template.py     (editable)  |\n"
        "      |   SSH    |                                            |\n"
        "      +--------->|  cc_manager.py  -> configures ADE remotely |\n"
        "                 |  Volume socx-sim-data -> /app/data         |\n"
        "                 |     |-- configured_ccs.json                |\n"
        "                 |     +-- response_template.json             |\n"
        "                 +--------------------------------------------+"
    )
    _code_shape(s, diagram, top=CONTENT_TOP, height=Emu(3700000), size=12)
    tf = _text_shape(s, top=Emu(5250000), height=Emu(950000))
    _para(tf, "Per-request destination matching: the response destinationIPs is always overwritten with the request network, so one global template serves many CCs simultaneously. State is file-based (no DB).", size=13, color=MUTED, first=True)
    s = prs.slides.add_slide(content)
    _set_title(s, "Architecture - key components")
    _remove_ph(s, CONTENT_LEFT)
    _remove_ph(s, CONTENT_RIGHT)
    rows = [
        ("main.py", "ASGI entrypoint; serves the app at / and /socx_sim/"),
        ("cloud_mock_server.py", "FastAPI app: _getRecommendation, UI + JSON APIs"),
        ("permanent_responses.py", "Pinned, exact rule sets for specific networks"),
        ("response_template.py", "Global, editable template (dst = request network)"),
        ("cc_manager.py", "Configures Cyber Controllers over SSH; CC registry"),
        ("static/ui.html", "The web UI served at /ui"),
        ("Docker + volume", "One-command run; socx-sim-data persists /app/data"),
    ]
    tbl = s.shapes.add_table(len(rows) + 1, 2, CONTENT_X, CONTENT_TOP, CONTENT_W, Emu(4600000)).table
    tbl.columns[0].width = Emu(3300000)
    tbl.columns[1].width = Emu(6266276)
    for c, text in enumerate(["Component", "Role"]):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = RADWARE_BLUE
        run = cell.text_frame.paragraphs[0].add_run(); run.text = text
        run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for r_i, (a, b) in enumerate(rows, start=1):
        for c_i, text in enumerate((a, b)):
            cell = tbl.cell(r_i, c_i)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run = cell.text_frame.paragraphs[0].add_run(); run.text = text
            run.font.size = Pt(12); run.font.color.rgb = INK
            if c_i == 0:
                run.font.name = "Consolas"; run.font.bold = True
    s = prs.slides.add_slide(content)
    _set_title(s, "The user interface  -  https://<host>:8080/ui")
    _fill_bullets(_ph(s, CONTENT_LEFT), [
        ("Tab 1 - Cyber Controller", 0, True),
        ("Enter CC host + SSH creds (sim address auto-detected)", 1, False),
        ("One click: set socx.*.cloud.hostname, import the TLS cert into ADE, restart ADE", 1, False),
        ("Test connection runs read-only preflight; Configure aborts before changing anything if a check fails", 1, False),
        ("Per-row Reset restores a CC to its original state", 1, False),
    ])
    _fill_bullets(_ph(s, CONTENT_RIGHT), [
        ("Tab 2 - Recommendation", 0, True),
        ("Edit the JSON returned for non-pinned requests", 1, False),
        ("destinationIPs always matches the request; other fields optional", 1, False),
        ("Preview shows the exact response a CC receives", 1, False),
        ("Tab 3 - Recommendations", 0, True),
        ("Browse pinned + seeded rule sets in a table", 1, False),
        ("View JSON, and Copy to Template to reuse one", 1, False),
    ])
    s = prs.slides.add_slide(content)
    _set_title(s, "CLI access to the stored files")
    _remove_ph(s, CONTENT_LEFT)
    _remove_ph(s, CONTENT_RIGHT)
    tf = _text_shape(s, top=CONTENT_TOP, height=Emu(700000))
    _para(tf, "Both state files live in the Docker volume socx-sim-data, mounted inside the container at /app/data.", size=15, color=MUTED, first=True)
    tf2 = _text_shape(s, top=Emu(2050000), height=Emu(430000))
    _para(tf2, "Registered Cyber Controllers  ->  configured_ccs.json", size=15, bold=True, color=RADWARE_BLUE, first=True)
    _code_shape(s, "docker exec socx-sim cat /app/data/configured_ccs.json | python3 -m json.tool", top=Emu(2470000), height=Emu(560000), size=12)
    tf3 = _text_shape(s, top=Emu(3230000), height=Emu(430000))
    _para(tf3, "All recommendations (editable template)  ->  response_template.json", size=15, bold=True, color=RADWARE_BLUE, first=True)
    _code_shape(s, "docker exec socx-sim cat /app/data/response_template.json | python3 -m json.tool", top=Emu(3650000), height=Emu(560000), size=12)
    tf4 = _text_shape(s, top=Emu(4450000), height=Emu(700000))
    _para(tf4, "Note: the pinned per-network recommendations are defined in code (permanent_responses.py), not in the data volume.", size=12, color=MUTED, first=True)
    s = prs.slides.add_slide(content)
    _set_title(s, "CLI access - via the host volume")
    _remove_ph(s, CONTENT_LEFT)
    _remove_ph(s, CONTENT_RIGHT)
    tf = _text_shape(s, top=CONTENT_TOP, height=Emu(560000))
    _para(tf, "If the container is not running, read the files straight from the volume on the host:", size=15, color=MUTED, first=True)
    host_cmds = (
        "# Find where Docker stores the volume on disk\n"
        "VOL=$(docker volume inspect socx-sim-data --format '{{ .Mountpoint }}')\n\n"
        "# Registered CCs\n"
        "sudo cat \"$VOL/configured_ccs.json\"    | python3 -m json.tool\n\n"
        "# All recommendations (template)\n"
        "sudo cat \"$VOL/response_template.json\" | python3 -m json.tool"
    )
    _code_shape(s, host_cmds, top=Emu(1950000), height=Emu(2450000), size=12)
    _code_shape(s, "# Copy a file out of the container to the current directory\ndocker cp socx-sim:/app/data/configured_ccs.json ./configured_ccs.json", top=Emu(4550000), height=Emu(900000), size=12)
    s = prs.slides.add_slide(content)
    _set_title(s, "Installation guide")
    _remove_ph(s, CONTENT_LEFT)
    _remove_ph(s, CONTENT_RIGHT)
    tf = _text_shape(s, top=CONTENT_TOP, height=Emu(1050000))
    _para(tf, "Canonical, always up to date - README on GitHub:", size=15, bold=True, color=RADWARE_BLUE, first=True)
    _para(tf, "https://github.com/rdwr-danaaz/cloudSimulator#deploy-to-a-new-machine-running-ade-step-by-step", size=13, color=INK)
    tf3 = _text_shape(s, top=Emu(2650000), height=Emu(450000))
    _para(tf3, "One-command install (run on the Linux host that will run the simulator):", size=15, bold=True, color=RADWARE_BLUE, first=True)
    _code_shape(s, "pip install -r deploy/requirements-deploy.txt\ncp deploy/install_config.example.json deploy/install_config.json   # add SSH creds\npython deploy/install.py", top=Emu(3150000), height=Emu(1250000), size=13)
    tf4 = _text_shape(s, top=Emu(4600000), height=Emu(600000))
    _para(tf4, "Nothing is installed manually on the Cyber Controller (ADE) - the UI configures it remotely over SSH.", size=12, color=MUTED, first=True)
    s = prs.slides.add_slide(cover)
    _ph(s, COVER_TITLE).text = "Thank you"
    _ph(s, COVER_LINE1).text = "SOC-X Recommendation Simulator"
    _ph(s, COVER_LINE2).text = "github.com/rdwr-danaaz/cloudSimulator"
    _ph(s, COVER_LINE3).text = "Web UI: https://<host>:8080/ui"
    prs.save(OUT)
    print("Wrote " + OUT)
if __name__ == "__main__":
    if "--force" in sys.argv:
        OUT = FINAL
    elif os.path.exists(FINAL):
        print("NOTE: PRESENTATION.pptx exists and is treated as the curated final deck.")
        print("      Writing to PRESENTATION_generated.pptx instead (use --force to overwrite).")
    build()
